import argparse
import os
from PIL import Image
import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation
from mdutils.mdutils import MdUtils
from mdutils import Html

def main():

    parser = argparse.ArgumentParser()
    parser.add_argument('ppt_name', type=str, help='add the name of the PowerPoint file(NOTE: the folder must be in the same directory as the prorgram file')
    args = parser.parse_args()
    
    pptx_name = args.ppt_name
    pptx_name_formatted = pptx_name.split('.')[0]

    prs = Presentation(pptx_name)

    path = '{}_converted'.format(pptx_name_formatted)
    if not os.path.exists(path):
        os.mkdir(path)
    images_folder = '{}_images'.format(pptx_name_formatted)
    images_path = os.path.join(path, images_folder)
    if not os.path.exists(images_path):
        os.mkdir(images_path)

    ppt_dict = {} #Keys: slide numbers, values: slide content
    texts = []
    slide_count = 0
    picture_count = 0
    for slide in prs.slides:
        texts = []
        slide_count += 1

        slide_parts = list(slide._part.related_parts.keys())
        for part in slide_parts:
            image_part = slide._part.related_parts[part]
            if type(image_part) == pptx.parts.image.ImagePart or pptx.opc.package.Part:
                file_startswith = image_part.blob[0:1]
                if file_startswith == b'\x89' or file_startswith == b'\xff' or file_startswith == b'\x47':
                    with open('{}/image{}_slide{}.png'.format(images_path, picture_count, slide_count), 'wb') as f:
                        f.write(image_part.blob)
                        picture_count += 1
                        
        for shape in slide.shapes:
            if shape.has_text_frame:
                if '\n' in shape.text:
                    splitted = shape.text.split('\n')
                    for word in splitted:
                        if word != '':
                            texts.append(word)
                elif shape.text == '':
                    continue
                else:
                    texts.append(shape.text)
            ppt_dict[slide_count] = texts

    ppt_content = ''
    for k,v in ppt_dict.items():
        ppt_content = ppt_content + ' - Slide number {}\n'.format(k)
        for a in v:
            ppt_content = ppt_content + '\t - {}\n'.format(a)

    mdFile = MdUtils(file_name='{}/{}'.format(path,path)) #second argument isn't path, it just shares the path name.
    mdFile.write(ppt_content)
    mdFile.create_md_file()


if __name__ == "__main__":
    main()